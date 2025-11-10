from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

# Create new presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# SLIDE 1: Baseline Burden
blank_layout = prs.slide_layouts[6]
slide1 = prs.slides.add_slide(blank_layout)

# Background
bg = slide1.background
bg.fill.solid()
bg.fill.fore_color.rgb = RGBColor(20, 30, 48)

# Title
title1 = slide1.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
tf = title1.text_frame
tf.text = "Abu Dhabi National Members – Health-Spend Burden (Thiqa)"
p = tf.paragraphs[0]
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER

# Left column text
left_box = slide1.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4.5), Inches(5))
tf = left_box.text_frame
tf.word_wrap = True

content = [
    ("Total Thiqa members (assumption): ", "1,000,000", False),
    ("Annual spend by Thiqa: ", "25 billion AED", True),
    ("", "", False),
    ("Cost distribution assumption:", "", True),
    ("", "", False),
    ("High-burden 20% of members → ", "200,000 people", False),
    ("High-burden spend: 80% of total = ", "20 billion AED", True),
    ("Average cost per person (high-burden): ", "100,000 AED/person/year", True),
    ("", "", False),
    ("Lower-burden 80% of members → ", "800,000 people", False),
    ("Lower-burden spend: 20% of total = ", "5 billion AED", True),
    ("Average cost per person (lower-burden): ", "6,250 AED/person/year", True),
]

for i, (label, value, highlight) in enumerate(content):
    if i > 0:
        tf.add_paragraph()
    p = tf.paragraphs[i]
    p.text = label + value
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0, 180, 216) if highlight else RGBColor(200, 200, 200)
    p.space_after = Pt(6)

# Right column - Add chart placeholder and callouts
chart_data = CategoryChartData()
chart_data.categories = ['High-burden\n200k people', 'Lower-burden\n800k people']
chart_data.add_series('Spend (AED billions)', (20, 5))

x, y, cx, cy = Inches(5.5), Inches(1.5), Inches(4), Inches(3.5)
chart = slide1.shapes.add_chart(
    XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data
).chart

# Callout boxes
box1 = slide1.shapes.add_textbox(Inches(6), Inches(5.2), Inches(1.5), Inches(0.6))
tf1 = box1.text_frame
tf1.text = "100k AED"
tf1.paragraphs[0].font.size = Pt(20)
tf1.paragraphs[0].font.bold = True
tf1.paragraphs[0].font.color.rgb = RGBColor(193, 18, 31)
tf1.paragraphs[0].alignment = PP_ALIGN.CENTER
box1.line.color.rgb = RGBColor(193, 18, 31)
box1.line.width = Pt(2)

box2 = slide1.shapes.add_textbox(Inches(8), Inches(5.2), Inches(1.5), Inches(0.6))
tf2 = box2.text_frame
tf2.text = "6.25k AED"
tf2.paragraphs[0].font.size = Pt(20)
tf2.paragraphs[0].font.bold = True
tf2.paragraphs[0].font.color.rgb = RGBColor(0, 150, 100)
tf2.paragraphs[0].alignment = PP_ALIGN.CENTER
box2.line.color.rgb = RGBColor(0, 150, 100)
box2.line.width = Pt(2)

# Footer
footer1 = slide1.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.5))
tf = footer1.text_frame
tf.text = "Note: Figures are based on assumption model. Actual Thiqa spend and membership to be verified."
tf.paragraphs[0].font.size = Pt(11)
tf.paragraphs[0].font.color.rgb = RGBColor(150, 150, 150)
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# SLIDE 2: Intervention Scenario
slide2 = prs.slides.add_slide(blank_layout)

# Background
bg2 = slide2.background
bg2.fill.solid()
bg2.fill.fore_color.rgb = RGBColor(20, 30, 48)

# Title
title2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
tf = title2.text_frame
tf.text = "Healthspan Extension Intervention: Saving Potential"
p = tf.paragraphs[0]
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER

# Left column
left_box2 = slide2.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4.5), Inches(5))
tf = left_box2.text_frame
tf.word_wrap = True

content2 = [
    "Target cohort: 5% of high-burden segment → 10,000 people",
    "",
    "Current average cost (high-burden): ~100,000 AED/person/year",
    "",
    "Lower-burden group average cost: ~6,250 AED/person/year",
    "",
    "Cost differential per person when moved:",
    "100,000 − 6,250 = 93,750 AED/person/year",
    "",
    "Total annual savings for cohort:",
    "10,000 × 93,750 AED = 937,500,000 AED",
    "(~0.94 billion AED/year)",
]

for i, text in enumerate(content2):
    if i > 0:
        tf.add_paragraph()
    p = tf.paragraphs[i]
    p.text = text
    p.font.size = Pt(14)
    highlight = any(x in text for x in ["93,750", "937,500,000", "0.94 billion"])
    p.font.color.rgb = RGBColor(0, 180, 216) if highlight else RGBColor(200, 200, 200)
    p.font.bold = highlight
    p.space_after = Pt(6)

# Right column - Visual
# "+1 Healthy Year" box
box_healthy = slide2.shapes.add_textbox(Inches(5.5), Inches(2), Inches(2), Inches(1))
tf = box_healthy.text_frame
tf.text = "+1 Healthy Year"
tf.paragraphs[0].font.size = Pt(18)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = RGBColor(0, 180, 216)
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
box_healthy.line.color.rgb = RGBColor(0, 180, 216)
box_healthy.line.width = Pt(3)

# Arrow (simulated with text)
arrow = slide2.shapes.add_textbox(Inches(7.6), Inches(2.3), Inches(0.5), Inches(0.5))
tf = arrow.text_frame
tf.text = "→"
tf.paragraphs[0].font.size = Pt(36)
tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

# "Low-cost bracket" box
box_low = slide2.shapes.add_textbox(Inches(8.2), Inches(2), Inches(1.5), Inches(1))
tf = box_low.text_frame
tf.text = "Low-Cost\nBracket"
tf.paragraphs[0].font.size = Pt(16)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = RGBColor(0, 150, 100)
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
box_low.line.color.rgb = RGBColor(0, 150, 100)
box_low.line.width = Pt(3)

# Highlight savings
savings_box = slide2.shapes.add_textbox(Inches(5.5), Inches(3.5), Inches(4), Inches(1))
tf = savings_box.text_frame
tf.text = "≈ 0.94 billion AED\nannual savings"
for p in tf.paragraphs:
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 180, 216)
    p.alignment = PP_ALIGN.CENTER

# Small note
note = slide2.shapes.add_textbox(Inches(5.5), Inches(4.7), Inches(4), Inches(0.8))
tf = note.text_frame
tf.text = "Based on shifting 1% of total members\ninto lower cost bracket via healthspan extension."
tf.paragraphs[0].font.size = Pt(11)
tf.paragraphs[0].font.color.rgb = RGBColor(180, 180, 180)
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# Footer
footer2 = slide2.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.8))
tf = footer2.text_frame
tf.text = "Key Takeaways:\n• Even a small cohort (1% of total) generates meaningful savings\n• Realistic model uses cost difference between brackets, not full elimination\n• Consider intervention costs, cohort variance, and multi-year effects for accuracy"
tf.paragraphs[0].font.size = Pt(11)
tf.paragraphs[0].font.color.rgb = RGBColor(150, 150, 150)

# Save
output_path = '/home/ubuntu/Daman-Hayat-AI-Project/presentation/Thiqa_Savings_Analysis.pptx'
prs.save(output_path)
print(f"✅ Thiqa savings slides created!")
print(f"📁 Saved to: {output_path}")
