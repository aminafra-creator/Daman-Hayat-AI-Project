#!/usr/bin/env python3.11
"""
Hayat AI - Final Investor Presentation
Apple + BCG Style: Minimal text, maximum impact, data-driven
3 Dynamite Slides + Supporting Slides
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

def add_title_slide(prs, title, subtitle=""):
    """Add title slide with Apple-style minimalism"""
    blank_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(14), Inches(2))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(72)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Subtitle
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(14), Inches(1))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_frame.paragraphs[0].font.size = Pt(28)
        subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)
        subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_blocks, background_color=None):
    """Add content slide with BCG-style structure"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # Background color
    if background_color:
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = background_color
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(15), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    
    # Content blocks
    num_blocks = len(content_blocks)
    block_width = 14 / num_blocks
    
    for i, block in enumerate(content_blocks):
        x_pos = 0.5 + (i * block_width)
        content_box = slide.shapes.add_textbox(Inches(x_pos + 0.2), Inches(1.5), Inches(block_width - 0.4), Inches(6.5))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        for item in block:
            p = content_frame.add_paragraph()
            p.text = item['text']
            p.font.size = Pt(item.get('size', 18))
            p.font.bold = item.get('bold', False)
            p.font.color.rgb = item.get('color', RGBColor(0, 0, 0))
            p.space_before = Pt(item.get('space_before', 6))
            p.space_after = Pt(item.get('space_after', 6))
    
    return slide

# SLIDE 1: TITLE SLIDE
add_title_slide(prs, "Hayat AI", "Digital Twin Healthspan Platform for UAE")

# SLIDE 2: THE $11 BILLION CRISIS (DYNAMITE SLIDE 1)
slide2_content = [
    [
        {'text': '🔴 TRIPLE THREAT', 'size': 32, 'bold': True, 'color': RGBColor(200, 0, 0)},
        {'text': '', 'size': 12},
        {'text': 'Cardiovascular Disease', 'size': 24, 'bold': True},
        {'text': '40% of deaths', 'size': 48, 'bold': True, 'color': RGBColor(200, 0, 0)},
        {'text': '80% preventable', 'size': 20, 'color': RGBColor(100, 100, 100)},
        {'text': '', 'size': 12},
        {'text': 'Diabetes', 'size': 24, 'bold': True},
        {'text': '20.7% prevalence', 'size': 48, 'bold': True, 'color': RGBColor(200, 0, 0)},
        {'text': '2nd highest globally', 'size': 20, 'color': RGBColor(100, 100, 100)},
        {'text': '', 'size': 12},
        {'text': 'Cancer', 'size': 24, 'bold': True},
        {'text': '12% of deaths', 'size': 48, 'bold': True, 'color': RGBColor(200, 0, 0)},
        {'text': '40-50% preventable', 'size': 20, 'color': RGBColor(100, 100, 100)},
    ],
    [
        {'text': '💰 ECONOMIC BURDEN', 'size': 32, 'bold': True, 'color': RGBColor(200, 100, 0)},
        {'text': '', 'size': 12},
        {'text': '$11 Billion', 'size': 64, 'bold': True, 'color': RGBColor(200, 100, 0)},
        {'text': 'Annual cost (AED 40B)', 'size': 20, 'color': RGBColor(100, 100, 100)},
        {'text': '', 'size': 12},
        {'text': '60-70%', 'size': 64, 'bold': True, 'color': RGBColor(200, 100, 0)},
        {'text': 'Preventable burden', 'size': 20, 'color': RGBColor(100, 100, 100)},
        {'text': '', 'size': 12},
        {'text': '16 Years', 'size': 64, 'bold': True, 'color': RGBColor(200, 100, 0)},
        {'text': 'Healthspan gap', 'size': 20, 'color': RGBColor(100, 100, 100)},
        {'text': '7th worst globally', 'size': 18, 'color': RGBColor(100, 100, 100)},
    ],
    [
        {'text': '🎯 AT RISK NOW', 'size': 32, 'bold': True, 'color': RGBColor(150, 0, 0)},
        {'text': '', 'size': 12},
        {'text': '213,000', 'size': 64, 'bold': True, 'color': RGBColor(150, 0, 0)},
        {'text': 'Thiqa members', 'size': 20, 'color': RGBColor(100, 100, 100)},
        {'text': 'Age 50+, chronic disease', 'size': 18, 'color': RGBColor(100, 100, 100)},
        {'text': '', 'size': 12},
        {'text': 'AED 20-30B', 'size': 48, 'bold': True, 'color': RGBColor(150, 0, 0)},
        {'text': 'Annual citizen healthcare', 'size': 20, 'color': RGBColor(100, 100, 100)},
        {'text': 'Abu Dhabi Government', 'size': 18, 'color': RGBColor(100, 100, 100)},
        {'text': '', 'size': 12},
        {'text': '27% of members', 'size': 28, 'bold': True},
        {'text': '= 60-70% of costs', 'size': 28, 'bold': True, 'color': RGBColor(150, 0, 0)},
    ]
]
add_content_slide(prs, "The $11 Billion Crisis", slide2_content)

# SLIDE 3: WHY DIGITAL TWIN AI WINS (DYNAMITE SLIDE 2)
slide3_content = [
    [
        {'text': '✅ PROVEN TECHNOLOGY', 'size': 28, 'bold': True, 'color': RGBColor(0, 150, 0)},
        {'text': '', 'size': 12},
        {'text': '93%', 'size': 56, 'bold': True, 'color': RGBColor(0, 150, 0)},
        {'text': 'Cancer detection accuracy', 'size': 18, 'color': RGBColor(100, 100, 100)},
        {'text': 'Nature, 373 citations', 'size': 16, 'color': RGBColor(100, 100, 100)},
        {'text': '', 'size': 8},
        {'text': '90%', 'size': 56, 'bold': True, 'color': RGBColor(0, 150, 0)},
        {'text': 'Cost reduction', 'size': 18, 'color': RGBColor(100, 100, 100)},
        {'text': 'Chronic disease mgmt', 'size': 16, 'color': RGBColor(100, 100, 100)},
        {'text': '', 'size': 8},
        {'text': '20+ Years', 'size': 48, 'bold': True, 'color': RGBColor(0, 150, 0)},
        {'text': 'Digital Twin research', 'size': 18, 'color': RGBColor(100, 100, 100)},
        {'text': 'Mount Sinai, 12M patients', 'size': 16, 'color': RGBColor(100, 100, 100)},
    ],
    [
        {'text': '💰 MARKET VALIDATION', 'size': 28, 'bold': True, 'color': RGBColor(0, 100, 200)},
        {'text': '', 'size': 12},
        {'text': 'Twin Health', 'size': 24, 'bold': True},
        {'text': '$950M', 'size': 56, 'bold': True, 'color': RGBColor(0, 100, 200)},
        {'text': 'Valuation (Aug 2025)', 'size': 18, 'color': RGBColor(100, 100, 100)},
        {'text': '$296M raised', 'size': 20, 'color': RGBColor(100, 100, 100)},
        {'text': 'Metabolic Digital Twin', 'size': 16, 'color': RGBColor(100, 100, 100)},
        {'text': '', 'size': 8},
        {'text': 'PrediSurge', 'size': 24, 'bold': True},
        {'text': '€6.5M', 'size': 48, 'bold': True, 'color': RGBColor(0, 100, 200)},
        {'text': 'Series A (Oct 2023)', 'size': 18, 'color': RGBColor(100, 100, 100)},
        {'text': 'Cardiovascular Digital Twin', 'size': 16, 'color': RGBColor(100, 100, 100)},
        {'text': '', 'size': 8},
        {'text': '→ Investors believe in', 'size': 18, 'bold': True},
        {'text': 'Digital Twin health tech', 'size': 18, 'bold': True, 'color': RGBColor(0, 100, 200)},
    ],
    [
        {'text': '🛡️ UAE-SPECIFIC MOAT', 'size': 28, 'bold': True, 'color': RGBColor(100, 0, 150)},
        {'text': '', 'size': 12},
        {'text': '1. Emirati Genetics', 'size': 20, 'bold': True},
        {'text': '800K genome database', 'size': 18, 'color': RGBColor(100, 100, 100)},
        {'text': '30-40% higher accuracy', 'size': 18, 'color': RGBColor(100, 0, 150)},
        {'text': '', 'size': 6},
        {'text': '2. Ramadan Protocols', 'size': 20, 'bold': True},
        {'text': 'Cultural adaptation', 'size': 18, 'color': RGBColor(100, 100, 100)},
        {'text': '40%+ adherence boost', 'size': 18, 'color': RGBColor(100, 0, 150)},
        {'text': '', 'size': 6},
        {'text': '3. Government Access', 'size': 20, 'bold': True},
        {'text': 'Daman/Thiqa partnership', 'size': 18, 'color': RGBColor(100, 100, 100)},
        {'text': '790K members', 'size': 18, 'color': RGBColor(100, 0, 150)},
        {'text': '', 'size': 6},
        {'text': '4. 6-Source Integration', 'size': 20, 'bold': True},
        {'text': 'vs competitors\' 1-2', 'size': 18, 'color': RGBColor(100, 0, 150)},
        {'text': '', 'size': 6},
        {'text': '→ Western competitors', 'size': 16, 'bold': True},
        {'text': 'can\'t replicate', 'size': 16, 'bold': True, 'color': RGBColor(100, 0, 150)},
    ]
]
add_content_slide(prs, "Why Digital Twin AI Wins", slide3_content)

# SLIDE 4: 7.6X RETURNS IN 5 YEARS (DYNAMITE SLIDE 3)
slide4_content = [
    [
        {'text': '💵 THE ASK', 'size': 32, 'bold': True, 'color': RGBColor(150, 100, 0)},
        {'text': '', 'size': 12},
        {'text': '$15M', 'size': 72, 'bold': True, 'color': RGBColor(150, 100, 0)},
        {'text': 'Seed raise', 'size': 24, 'color': RGBColor(100, 100, 100)},
        {'text': '', 'size': 12},
        {'text': '$35M pre-money', 'size': 24},
        {'text': '$50M post-money', 'size': 24},
        {'text': '', 'size': 12},
        {'text': 'Use of Funds:', 'size': 22, 'bold': True},
        {'text': '• Team: $4M (27%)', 'size': 18},
        {'text': '• Product: $6M (40%)', 'size': 18},
        {'text': '• Clinical: $3M (20%)', 'size': 18},
        {'text': '• Operations: $2M (13%)', 'size': 18},
    ],
    [
        {'text': '📈 5-YEAR PROJECTIONS', 'size': 32, 'bold': True, 'color': RGBColor(100, 0, 150)},
        {'text': '', 'size': 12},
        {'text': 'Year 2 (Series A)', 'size': 22, 'bold': True},
        {'text': '5,925 members', 'size': 20},
        {'text': '$7.1M ARR', 'size': 28, 'bold': True, 'color': RGBColor(100, 0, 150)},
        {'text': '', 'size': 8},
        {'text': 'Year 3 (Break-even)', 'size': 22, 'bold': True},
        {'text': '14,036 members', 'size': 20},
        {'text': '$16.8M ARR', 'size': 28, 'bold': True, 'color': RGBColor(100, 0, 150)},
        {'text': '', 'size': 8},
        {'text': 'Year 5 (Exit Ready)', 'size': 22, 'bold': True},
        {'text': '38,191 members', 'size': 20},
        {'text': '$45.8M ARR', 'size': 28, 'bold': True, 'color': RGBColor(100, 0, 150)},
        {'text': '49% EBITDA margin', 'size': 18, 'color': RGBColor(100, 100, 100)},
        {'text': '', 'size': 8},
        {'text': '18% penetration', 'size': 20},
        {'text': '15% churn factored in', 'size': 18, 'color': RGBColor(100, 100, 100)},
    ],
    [
        {'text': '🎯 RETURNS', 'size': 32, 'bold': True, 'color': RGBColor(200, 150, 0)},
        {'text': '', 'size': 12},
        {'text': '7.6x', 'size': 96, 'bold': True, 'color': RGBColor(200, 150, 0)},
        {'text': 'Returns in 5 years', 'size': 24, 'color': RGBColor(100, 100, 100)},
        {'text': '', 'size': 12},
        {'text': '50% IRR', 'size': 48, 'bold': True, 'color': RGBColor(200, 150, 0)},
        {'text': 'Internal rate of return', 'size': 20, 'color': RGBColor(100, 100, 100)},
        {'text': '', 'size': 12},
        {'text': '$380M', 'size': 56, 'bold': True, 'color': RGBColor(200, 150, 0)},
        {'text': 'Valuation (10x ARR)', 'size': 20, 'color': RGBColor(100, 100, 100)},
        {'text': '', 'size': 12},
        {'text': 'Upside: 10.5x', 'size': 28, 'bold': True},
        {'text': 'at 25% penetration (Y7)', 'size': 18, 'color': RGBColor(100, 100, 100)},
        {'text': 'Twin Health: 30%+', 'size': 16, 'color': RGBColor(100, 100, 100)},
    ]
]
add_content_slide(prs, "7.6x Returns in 5 Years", slide4_content)

# SLIDE 5: EXECUTION ROADMAP
slide5_content = [
    [
        {'text': 'SEED (M0-18)', 'size': 28, 'bold': True, 'color': RGBColor(0, 100, 200)},
        {'text': '', 'size': 8},
        {'text': 'M0-6: Build', 'size': 20, 'bold': True},
        {'text': '• Daman MOU signed', 'size': 16},
        {'text': '• Team assembled', 'size': 16},
        {'text': '• Platform MVP', 'size': 16},
        {'text': '', 'size': 6},
        {'text': 'M6-12: Validate', 'size': 20, 'bold': True},
        {'text': '• 500-person pilot', 'size': 16},
        {'text': '• Thiqa members, 50+', 'size': 16},
        {'text': '• Health Score +10', 'size': 16},
        {'text': '', 'size': 6},
        {'text': 'M12-18: Prove', 'size': 20, 'bold': True},
        {'text': '• 75%+ retention', 'size': 16},
        {'text': '• Daman LOI (10K)', 'size': 16},
        {'text': '• $600K ARR', 'size': 16},
    ],
    [
        {'text': 'SERIES A (Y2)', 'size': 28, 'bold': True, 'color': RGBColor(0, 150, 0)},
        {'text': '', 'size': 8},
        {'text': 'Raise: $15-20M', 'size': 22, 'bold': True},
        {'text': 'Pre-money: $50-60M', 'size': 18},
        {'text': '', 'size': 8},
        {'text': 'Metrics:', 'size': 20, 'bold': True},
        {'text': '• 5,925 members', 'size': 16},
        {'text': '• $7.1M ARR', 'size': 16},
        {'text': '• 85% retention', 'size': 16},
        {'text': '• Daman LOI executed', 'size': 16},
        {'text': '', 'size': 8},
        {'text': 'Goal:', 'size': 20, 'bold': True},
        {'text': '• Scale to 14K (Y3)', 'size': 16},
        {'text': '• Break-even', 'size': 16},
    ],
    [
        {'text': 'SCALE (Y3-5)', 'size': 28, 'bold': True, 'color': RGBColor(150, 0, 150)},
        {'text': '', 'size': 8},
        {'text': 'Year 3:', 'size': 20, 'bold': True},
        {'text': '• 14,036 members', 'size': 16},
        {'text': '• $16.8M ARR', 'size': 16},
        {'text': '• Break-even', 'size': 16},
        {'text': '', 'size': 6},
        {'text': 'Year 4:', 'size': 20, 'bold': True},
        {'text': '• 24,931 members', 'size': 16},
        {'text': '• $29.9M ARR', 'size': 16},
        {'text': '• 37% EBITDA', 'size': 16},
        {'text': '', 'size': 6},
        {'text': 'Year 5:', 'size': 20, 'bold': True},
        {'text': '• 38,191 members', 'size': 16},
        {'text': '• $45.8M ARR', 'size': 16},
        {'text': '• 49% EBITDA', 'size': 16},
        {'text': '• $380M valuation', 'size': 16, 'bold': True, 'color': RGBColor(150, 0, 150)},
    ]
]
add_content_slide(prs, "Execution Roadmap: Seed → Series A → Exit", slide5_content)

# SLIDE 6: WHY THIQA/DAMAN?
slide6_content = [
    [
        {'text': '🏛️ THIQA STRUCTURE', 'size': 28, 'bold': True, 'color': RGBColor(0, 100, 150)},
        {'text': '', 'size': 8},
        {'text': 'What is Thiqa?', 'size': 22, 'bold': True},
        {'text': '• Abu Dhabi Gov program', 'size': 16},
        {'text': '• 100% coverage for UAE', 'size': 16},
        {'text': '  nationals', 'size': 16},
        {'text': '• Managed by Daman', 'size': 16},
        {'text': '• Fully gov-funded', 'size': 16},
        {'text': '', 'size': 6},
        {'text': 'Members: 790,000', 'size': 20, 'bold': True},
        {'text': 'Spend: AED 20-30B', 'size': 20, 'bold': True},
        {'text': '', 'size': 6},
        {'text': 'Source:', 'size': 14, 'color': RGBColor(100, 100, 100)},
        {'text': 'Daman 10th Anniversary', 'size': 14, 'color': RGBColor(100, 100, 100)},
        {'text': '(2018)', 'size': 14, 'color': RGBColor(100, 100, 100)},
    ],
    [
        {'text': '🎯 WHY TARGET 50+?', 'size': 28, 'bold': True, 'color': RGBColor(150, 0, 0)},
        {'text': '', 'size': 8},
        {'text': 'Target: 213,000', 'size': 24, 'bold': True, 'color': RGBColor(150, 0, 0)},
        {'text': '(27% of Thiqa)', 'size': 18, 'color': RGBColor(100, 100, 100)},
        {'text': '', 'size': 8},
        {'text': '= 60-70%', 'size': 48, 'bold': True, 'color': RGBColor(150, 0, 0)},
        {'text': 'of healthcare costs', 'size': 20, 'color': RGBColor(100, 100, 100)},
        {'text': '', 'size': 8},
        {'text': 'Per-capita spend:', 'size': 20, 'bold': True},
        {'text': 'AED 61-91K/year', 'size': 22, 'bold': True, 'color': RGBColor(150, 0, 0)},
        {'text': '($16,600-24,900)', 'size': 18, 'color': RGBColor(100, 100, 100)},
        {'text': '', 'size': 6},
        {'text': '→ Highest-cost,', 'size': 18, 'bold': True},
        {'text': 'highest-value segment', 'size': 18, 'bold': True, 'color': RGBColor(150, 0, 0)},
    ],
    [
        {'text': '💰 ROI FOR DAMAN', 'size': 28, 'bold': True, 'color': RGBColor(0, 150, 0)},
        {'text': '', 'size': 8},
        {'text': 'Annual cost: $16-25K', 'size': 18},
        {'text': 'Preventable (70%):', 'size': 18},
        {'text': '$11-17K', 'size': 20, 'bold': True},
        {'text': '', 'size': 6},
        {'text': 'Hayat AI reduces 20%:', 'size': 18},
        {'text': '$2-5K savings', 'size': 24, 'bold': True, 'color': RGBColor(0, 150, 0)},
        {'text': '', 'size': 6},
        {'text': 'Hayat AI cost:', 'size': 18},
        {'text': '$1,200/year', 'size': 24, 'bold': True},
        {'text': '', 'size': 6},
        {'text': 'ROI: 1.9-4.4x', 'size': 32, 'bold': True, 'color': RGBColor(0, 150, 0)},
        {'text': '', 'size': 6},
        {'text': 'At 38K members:', 'size': 18},
        {'text': '$112-403M', 'size': 28, 'bold': True, 'color': RGBColor(0, 150, 0)},
        {'text': 'annual net savings', 'size': 18, 'color': RGBColor(100, 100, 100)},
    ]
]
add_content_slide(prs, "Why Thiqa via Daman? Government-Backed, High-Value", slide6_content)

# SLIDE 7: COMPETITIVE ADVANTAGES
slide7_content = [
    [
        {'text': 'UAE-SPECIFIC MOAT', 'size': 28, 'bold': True, 'color': RGBColor(100, 0, 150)},
        {'text': '', 'size': 8},
        {'text': '1. Emirati Genetics', 'size': 22, 'bold': True},
        {'text': '800K genome database', 'size': 16},
        {'text': '30-40% accuracy boost', 'size': 16, 'color': RGBColor(100, 0, 150)},
        {'text': '', 'size': 6},
        {'text': '2. Ramadan Protocols', 'size': 22, 'bold': True},
        {'text': 'Cultural adaptation', 'size': 16},
        {'text': '40%+ adherence boost', 'size': 16, 'color': RGBColor(100, 0, 150)},
        {'text': '', 'size': 6},
        {'text': '3. Arabic Language', 'size': 22, 'bold': True},
        {'text': 'Native interface', 'size': 16},
        {'text': 'Broader accessibility', 'size': 16, 'color': RGBColor(100, 0, 150)},
    ],
    [
        {'text': 'PARTNERSHIP MOAT', 'size': 28, 'bold': True, 'color': RGBColor(0, 100, 150)},
        {'text': '', 'size': 8},
        {'text': '4. Government Access', 'size': 22, 'bold': True},
        {'text': 'Daman/Thiqa partnership', 'size': 16},
        {'text': '790K members', 'size': 16, 'color': RGBColor(0, 100, 150)},
        {'text': '', 'size': 6},
        {'text': '5. First-Mover', 'size': 22, 'bold': True},
        {'text': 'UAE Digital Twin market', 'size': 16},
        {'text': 'Network effects', 'size': 16, 'color': RGBColor(0, 100, 150)},
        {'text': '', 'size': 6},
        {'text': '6. Regulatory', 'size': 22, 'bold': True},
        {'text': 'Streamlined approval', 'size': 16},
        {'text': 'Gov-run Thiqa', 'size': 16, 'color': RGBColor(0, 100, 150)},
    ],
    [
        {'text': 'TECHNICAL MOAT', 'size': 28, 'bold': True, 'color': RGBColor(150, 100, 0)},
        {'text': '', 'size': 8},
        {'text': '7. 6-Source Integration', 'size': 22, 'bold': True},
        {'text': 'vs competitors\' 1-2', 'size': 16},
        {'text': 'Comprehensive data', 'size': 16, 'color': RGBColor(150, 100, 0)},
        {'text': '', 'size': 6},
        {'text': '8. Blueprint-Optimized', 'size': 22, 'bold': True},
        {'text': 'Evidence-based', 'size': 16},
        {'text': 'Health Score V2.0', 'size': 16, 'color': RGBColor(150, 100, 0)},
        {'text': '', 'size': 6},
        {'text': '→ Western competitors', 'size': 18, 'bold': True},
        {'text': 'face 3-5 year', 'size': 18},
        {'text': 'market entry barriers', 'size': 18, 'bold': True, 'color': RGBColor(150, 100, 0)},
    ]
]
add_content_slide(prs, "8 Defensible Competitive Advantages", slide7_content)

# SLIDE 8: THANK YOU
add_title_slide(prs, "Thank You", "Questions?")

# Save presentation
output_file = "/home/ubuntu/Daman-Hayat-AI-Project/Hayat_AI_Final_Presentation.pptx"
prs.save(output_file)
print(f"✅ Presentation created: {output_file}")
print(f"📊 Total slides: {len(prs.slides)}")
print(f"🎯 Focus: First 3 dynamite slides (Slides 2-4)")
print(f"🎨 Style: Apple minimalism + BCG rigor")
