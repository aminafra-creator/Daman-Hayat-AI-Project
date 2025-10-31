from pptx import Presentation
import json

pptx_path = "/home/ubuntu/Daman-Hayat-AI-Project/presentation/Daman-Hayat_AI_Transforming_Prevention_into_Profit.pptx"
prs = Presentation(pptx_path)

print(f"Total slides: {len(prs.slides)}")
print(f"Slide dimensions: {prs.slide_width} x {prs.slide_height}")
print("\n" + "="*80 + "\n")

for idx, slide in enumerate(prs.slides, 1):
    print(f"SLIDE {idx}")
    print("-" * 40)
    
    # Extract all text from the slide
    text_content = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            text_content.append(shape.text.strip())
    
    if text_content:
        for text in text_content:
            print(f"  {text}")
    else:
        print("  [No text content]")
    
    print("\n")

print("="*80)
print(f"\nAnalysis complete. Total slides: {len(prs.slides)}")
