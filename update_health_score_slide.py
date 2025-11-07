from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Load the presentation
prs = Presentation('/home/ubuntu/Daman-Hayat-AI-Project/presentation/Daman-Hayat_AI_Opportunity_HealthScore.pptx')

# Get the last slide (Health Score slide)
slide = prs.slides[-1]

# Update subtitle to emphasize Digital Twin
for shape in slide.shapes:
    if shape.has_text_frame:
        text = shape.text_frame.text
        if "One Comprehensive Score" in text:
            shape.text_frame.clear()
            p = shape.text_frame.paragraphs[0]
            p.text = "Powered by Your Personal Digital Twin • Predicts Your Health 12-18 Months Ahead"
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(0, 180, 216)  # Cyan to emphasize
            p.alignment = PP_ALIGN.CENTER

# Update bottom text to emphasize Digital Twin
for shape in slide.shapes:
    if shape.has_text_frame:
        text = shape.text_frame.text
        if "Based on Bryan Johnson" in text:
            shape.text_frame.clear()
            p = shape.text_frame.paragraphs[0]
            p.text = "Your Digital Twin learns from your unique biology • Blueprint-optimized weighting • UAE-customized norms"
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(150, 150, 150)
            p.alignment = PP_ALIGN.CENTER

# Save presentation
output_path = '/home/ubuntu/Daman-Hayat-AI-Project/presentation/Daman-Hayat_AI_Opportunity_HealthScore_v2.pptx'
prs.save(output_path)
print(f"✅ Health Score slide updated with Digital Twin emphasis!")
print(f"📁 Saved to: {output_path}")
