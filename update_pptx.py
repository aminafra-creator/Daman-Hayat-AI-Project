from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import re

# Load presentation
pptx_path = "/home/ubuntu/Daman-Hayat-AI-Project/presentation/Daman-Hayat_AI_Transforming_Prevention_into_Profit.pptx"
prs = Presentation(pptx_path)

print(f"Processing {len(prs.slides)} slides...")

# Currency conversion mapping (AED to USD, rounded)
currency_map = {
    'AED 55M': '$15M',
    'AED 19B': '$5.2B',
    'AED 60-75B': '$16-20B',
    'AED 5-8B': '$1.4-2.2B',
    'AED 750M-1.2B': '$200-330M',
    'AED 40B': '$11B',
    'AED 21B': '$5.7B',
    'AED 11B': '$3.0B',
    'AED 6.83B': '$1.9B',
    'AED 2.07B': '$560M',
    'AED 7,000': '$1,900',
    'AED 30,000': '$8,200',
    'AED 25,000': '$6,800',
    'AED 60,000': '$16,000',
    'AED 40,000': '$11,000',
    'AED 192,000': '$52,000',
    'AED 184,000': '$50,000',
    'AED 68B': '$19B',
    'AED 8,250': '$2,200',
    'AED 200-500': '$55-140',
    'AED 1,000-5,000': '$270-1,400',
    'AED 100': '$30',
    'AED 367-734': '$100-200',
}

slides_updated = []

# Process each slide
for idx, slide in enumerate(prs.slides, 1):
    slide_modified = False
    
    # Process all text in shapes
    for shape in slide.shapes:
        if hasattr(shape, "text_frame"):
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    original_text = run.text
                    new_text = original_text
                    
                    # Apply currency conversions
                    for aed, usd in currency_map.items():
                        if aed in new_text:
                            new_text = new_text.replace(aed, usd)
                            slide_modified = True
                    
                    # Update text if changed
                    if new_text != original_text:
                        run.text = new_text
    
    if slide_modified:
        slides_updated.append(idx)
        print(f"  ✓ Slide {idx}: Currency converted to USD")

print(f"\nCurrency conversion complete!")
print(f"Updated slides: {slides_updated}")
print(f"Total slides modified: {len(slides_updated)}")

# Save updated presentation
output_path = "/home/ubuntu/Daman-Hayat-AI-Project/presentation/Daman-Hayat_AI_USD_Updated.pptx"
prs.save(output_path)
print(f"\nSaved to: {output_path}")
