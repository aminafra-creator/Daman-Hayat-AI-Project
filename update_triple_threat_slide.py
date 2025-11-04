#!/usr/bin/env python3
"""
Update Triple Threat slide with:
1. Cancer prevention factors (40-50% preventable)
2. Diabetes clarification (24-25% in nationals)
3. Morbidity burden (60%+ of chronic disease burden)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# File paths
INPUT_FILE = "/home/ubuntu/Daman-Hayat-AI-Project/presentation/Daman-Hayat_AI_Opportunity_TripleThreat.pptx"
OUTPUT_FILE = "/home/ubuntu/Daman-Hayat-AI-Project/presentation/Daman-Hayat_AI_Opportunity_Updated.pptx"

# Color palette
COLOR_RED = RGBColor(193, 18, 31)
COLOR_ORANGE = RGBColor(255, 165, 0)
COLOR_PURPLE = RGBColor(114, 9, 183)
COLOR_TEAL = RGBColor(5, 102, 141)
COLOR_DARK_GRAY = RGBColor(44, 62, 80)
COLOR_LIGHT_GRAY = RGBColor(240, 240, 240)

def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=COLOR_DARK_GRAY, align=PP_ALIGN.CENTER):
    """Helper function to add text box with formatting"""
    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    text_frame = textbox.text_frame
    text_frame.text = text
    text_frame.word_wrap = True
    
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = align
    
    run = paragraph.runs[0]
    run.font.name = "Noto Sans"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    
    return textbox

def update_triple_threat_slide(prs):
    """Update the Triple Threat slide (last slide)"""
    
    # Get the last slide (Triple Threat)
    slide = prs.slides[-1]
    
    print(f"Updating slide with {len(slide.shapes)} shapes...")
    
    # Column dimensions (matching original)
    col_width = 2.8
    col_height = 5.0
    col_top = 1.5
    col_spacing = 0.2
    
    col1_left = 0.5
    col2_left = col1_left + col_width + col_spacing
    col3_left = col2_left + col_width + col_spacing
    
    # Update Diabetes column - add clarification
    # Find and update the "20.7%" text to include nationals clarification
    y_offset = col_top + 1.4  # Position after "20.7%"
    add_text_box(slide, col2_left, y_offset, col_width, 0.3,
                 "(24-25% in nationals)", font_size=12, color=COLOR_ORANGE, align=PP_ALIGN.CENTER)
    
    # Update Cancer column - add prevention statistic
    y_offset = col_top + 4.9  # After "10 years early detection"
    add_text_box(slide, col3_left, y_offset, col_width, 0.5,
                 "40-50%", font_size=36, bold=True, color=COLOR_TEAL)
    add_text_box(slide, col3_left, y_offset + 0.5, col_width, 0.3,
                 "preventable (lifestyle)", font_size=14, color=COLOR_TEAL)
    
    # Update Cancer detection methods to include prevention
    add_text_box(slide, col3_left, col_top + col_height - 0.5, col_width, 0.5,
                 "Detect: AI Biopsy, Genome | Prevent: Weight, Diet, Exercise", 
                 font_size=10, color=COLOR_DARK_GRAY)
    
    # Update bottom section to include morbidity
    bottom_top = col_top + col_height + 0.3
    
    # Update combined impact text
    add_text_box(slide, 0.5, bottom_top + 0.1, 9, 0.3,
                 "COMBINED IMPACT: 52% of deaths + 60%+ of chronic disease burden + $11B annual cost",
                 font_size=19, bold=True, color=COLOR_RED)
    
    return slide

def main():
    """Main function to update Triple Threat slide"""
    print("Loading presentation...")
    prs = Presentation(INPUT_FILE)
    
    print(f"Presentation has {len(prs.slides)} slides")
    
    print("Updating Triple Threat slide...")
    update_triple_threat_slide(prs)
    
    print(f"Saving to {OUTPUT_FILE}...")
    prs.save(OUTPUT_FILE)
    
    print("✅ Triple Threat slide updated successfully!")
    print(f"📁 Output file: {OUTPUT_FILE}")
    print("\nUpdates applied:")
    print("  ✅ Diabetes: Added '(24-25% in nationals)' clarification")
    print("  ✅ Cancer: Added '40-50% preventable (lifestyle)' statistic")
    print("  ✅ Cancer: Updated detection methods to include prevention")
    print("  ✅ Bottom: Added '60%+ of chronic disease burden' to impact")

if __name__ == "__main__":
    main()
