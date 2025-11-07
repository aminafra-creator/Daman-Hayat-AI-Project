#!/usr/bin/env python3.11
"""
Create 5 Dynamite Slides PowerPoint Presentation
Hayat AI Seed-Stage Investor Pitch
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

def create_presentation():
    """Create the 5 Dynamite Slides presentation"""
    
    # Create presentation with 16:9 aspect ratio
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # Define color palette
    colors = {
        'navy': RGBColor(26, 32, 44),      # Dark navy background
        'white': RGBColor(255, 255, 255),   # White text
        'light_gray': RGBColor(203, 213, 224),  # Light gray
        'red': RGBColor(255, 0, 0),         # Crisis/urgency
        'orange': RGBColor(204, 85, 0),     # Warning
        'gold': RGBColor(255, 215, 0),      # Opportunity
        'green': RGBColor(80, 200, 120),    # Success/growth
        'blue': RGBColor(0, 102, 204),      # Technology/trust
        'blue_gray': RGBColor(74, 85, 104), # Academic
    }
    
    # SLIDE 1: The $11 Billion Crisis
    print("Creating Slide 1: The $11 Billion Crisis...")
    slide1 = create_slide1(prs, colors)
    
    # SLIDE 2: Why Digital Twin AI Wins
    print("Creating Slide 2: Why Digital Twin AI Wins...")
    slide2 = create_slide2(prs, colors)
    
    # SLIDE 3: Seed → Series A → Dominance
    print("Creating Slide 3: Seed → Series A → Dominance...")
    slide3 = create_slide3(prs, colors)
    
    # SLIDE 4: UAE Healthspan Crisis Graph
    print("Creating Slide 4: UAE Healthspan Crisis Graph...")
    slide4 = create_slide4(prs, colors)
    
    # SLIDE 5: Digital Twin Market Validation
    print("Creating Slide 5: Digital Twin Market Validation...")
    slide5 = create_slide5(prs, colors)
    
    # Save presentation
    output_path = '/home/ubuntu/Daman-Hayat-AI-Project/Hayat_AI_5_Dynamite_Slides.pptx'
    prs.save(output_path)
    print(f"\n✅ Presentation saved: {output_path}")
    
    return output_path


def create_slide1(prs, colors):
    """Slide 1: The $11 Billion Crisis"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = colors['navy']
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "The $11 Billion Crisis"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = colors['white']
    title_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(0.4))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "UAE's Triple Threat: CVD, Diabetes, Cancer"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(32)
    subtitle_para.font.color.rgb = colors['light_gray']
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Three columns: Triple Threat, Economic Impact, Opportunity
    col_width = 2.8
    col_height = 3.0
    y_start = 1.8
    
    # Column 1: Triple Threat
    col1_x = 0.5
    col1_box = slide.shapes.add_shape(1, Inches(col1_x), Inches(y_start), Inches(col_width), Inches(col_height))
    col1_box.fill.solid()
    col1_box.fill.fore_color.rgb = RGBColor(40, 50, 70)
    col1_box.line.color.rgb = colors['red']
    col1_box.line.width = Pt(2)
    
    col1_text = col1_box.text_frame
    col1_text.word_wrap = True
    col1_text.margin_left = Inches(0.2)
    col1_text.margin_right = Inches(0.2)
    col1_text.margin_top = Inches(0.2)
    
    # Add Triple Threat content
    p1 = col1_text.paragraphs[0]
    p1.text = "🫀 TRIPLE THREAT"
    p1.font.size = Pt(28)
    p1.font.bold = True
    p1.font.color.rgb = colors['red']
    p1.alignment = PP_ALIGN.CENTER
    
    content1 = [
        "",
        "CVD: 40% of deaths",
        "2x cancer rate",
        "80% preventable",
        "",
        "Diabetes: 20.7% prevalence",
        "24-25% Emiratis",
        "90% preventable",
        "",
        "Cancer: 12% of deaths",
        "5,526 cases/year",
        "40-50% preventable",
        "",
        "Combined: 52% of deaths",
        "60%+ chronic burden"
    ]
    
    for line in content1:
        p = col1_text.add_paragraph()
        p.text = line
        p.font.size = Pt(16)
        p.font.color.rgb = colors['white']
        p.alignment = PP_ALIGN.LEFT
    
    # Column 2: Economic Impact
    col2_x = 3.6
    col2_box = slide.shapes.add_shape(1, Inches(col2_x), Inches(y_start), Inches(col_width), Inches(col_height))
    col2_box.fill.solid()
    col2_box.fill.fore_color.rgb = RGBColor(40, 50, 70)
    col2_box.line.color.rgb = colors['orange']
    col2_box.line.width = Pt(2)
    
    col2_text = col2_box.text_frame
    col2_text.word_wrap = True
    col2_text.margin_left = Inches(0.2)
    col2_text.margin_right = Inches(0.2)
    col2_text.margin_top = Inches(0.2)
    
    p2 = col2_text.paragraphs[0]
    p2.text = "💰 ECONOMIC IMPACT"
    p2.font.size = Pt(28)
    p2.font.bold = True
    p2.font.color.rgb = colors['orange']
    p2.alignment = PP_ALIGN.CENTER
    
    content2 = [
        "",
        "$11 Billion Annual Burden",
        "",
        "$5.2B Healthcare Costs",
        "$5.7B Lost Productivity",
        "",
        "16-Year Healthspan Gap",
        "7th worst globally",
        "",
        "Widened 45% Since 2000",
        "(Getting worse, not better)",
        "",
        "500,000+ Members",
        "Daman/Hayat at risk NOW"
    ]
    
    for line in content2:
        p = col2_text.add_paragraph()
        p.text = line
        if "$11" in line or "16-Year" in line:
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = colors['gold']
        else:
            p.font.size = Pt(16)
            p.font.color.rgb = colors['white']
        p.alignment = PP_ALIGN.LEFT
    
    # Column 3: Opportunity
    col3_x = 6.7
    col3_box = slide.shapes.add_shape(1, Inches(col3_x), Inches(y_start), Inches(col_width), Inches(col_height))
    col3_box.fill.solid()
    col3_box.fill.fore_color.rgb = RGBColor(40, 50, 70)
    col3_box.line.color.rgb = colors['green']
    col3_box.line.width = Pt(2)
    
    col3_text = col3_box.text_frame
    col3_text.word_wrap = True
    col3_text.margin_left = Inches(0.2)
    col3_text.margin_right = Inches(0.2)
    col3_text.margin_top = Inches(0.2)
    
    p3 = col3_text.paragraphs[0]
    p3.text = "🎯 OPPORTUNITY"
    p3.font.size = Pt(28)
    p3.font.bold = True
    p3.font.color.rgb = colors['green']
    p3.alignment = PP_ALIGN.CENTER
    
    content3 = [
        "",
        "60%+ Preventable",
        "",
        "$6.6B+ Savings Potential",
        "(60% of $11B burden)",
        "",
        "+1 Year Healthspan =",
        "+$50,000 Value/Person",
        "",
        "$25 Billion TAM",
        "(500K members × $50K)",
        "",
        "Digital Twin AI",
        "Can Reverse This Crisis"
    ]
    
    for line in content3:
        p = col3_text.add_paragraph()
        p.text = line
        if "$6.6B" in line or "$25 Billion" in line:
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = colors['gold']
        else:
            p.font.size = Pt(16)
            p.font.color.rgb = colors['white']
        p.alignment = PP_ALIGN.LEFT
    
    return slide


def create_slide2(prs, colors):
    """Slide 2: Why Digital Twin AI Wins in the UAE"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = colors['navy']
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Why Digital Twin AI Wins in the UAE"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = colors['white']
    title_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.95), Inches(9), Inches(0.4))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Proven Technology + UAE-Specific Moat"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.color.rgb = colors['light_gray']
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Proven Technology section (top)
    tech_box = slide.shapes.add_shape(1, Inches(0.5), Inches(1.5), Inches(9), Inches(1.0))
    tech_box.fill.solid()
    tech_box.fill.fore_color.rgb = RGBColor(40, 50, 70)
    tech_box.line.color.rgb = colors['green']
    tech_box.line.width = Pt(2)
    
    tech_text = tech_box.text_frame
    tech_text.word_wrap = True
    tech_text.margin_left = Inches(0.2)
    tech_text.margin_top = Inches(0.1)
    
    pt = tech_text.paragraphs[0]
    pt.text = "✅ PROVEN TECHNOLOGY FOUNDATION"
    pt.font.size = Pt(24)
    pt.font.bold = True
    pt.font.color.rgb = colors['green']
    pt.alignment = PP_ALIGN.CENTER
    
    tech_content = [
        "",
        "Digital Twin: 20+ years (Boeing, GE) | AI Blood: 93% cancer detection (Nature 2024) | Retinal: 90% cost reduction ($30 vs $270-1,400) | Google Health Partnership validates approach"
    ]
    
    for line in tech_content:
        p = tech_text.add_paragraph()
        p.text = line
        p.font.size = Pt(16)
        p.font.color.rgb = colors['white']
        p.alignment = PP_ALIGN.CENTER
    
    # Two columns: Multi-Modal Moat + UAE Advantages
    col_y = 2.7
    col_height = 2.3
    
    # Left: Multi-Modal Moat
    moat_box = slide.shapes.add_shape(1, Inches(0.5), Inches(col_y), Inches(4.4), Inches(col_height))
    moat_box.fill.solid()
    moat_box.fill.fore_color.rgb = RGBColor(40, 50, 70)
    moat_box.line.color.rgb = colors['blue']
    moat_box.line.width = Pt(2)
    
    moat_text = moat_box.text_frame
    moat_text.word_wrap = True
    moat_text.margin_left = Inches(0.2)
    moat_text.margin_right = Inches(0.2)
    moat_text.margin_top = Inches(0.1)
    
    pm = moat_text.paragraphs[0]
    pm.text = "🛡️ 6-SOURCE MOAT"
    pm.font.size = Pt(24)
    pm.font.bold = True
    pm.font.color.rgb = colors['blue']
    pm.alignment = PP_ALIGN.CENTER
    
    moat_content = [
        "",
        "1. Blood biomarkers ($55-140)",
        "2. Genome (800K UAE genomes)",
        "3. Medical imaging ($270-1,400)",
        "4. AI retinal screening ($30)",
        "5. Wearables (continuous)",
        "6. DEXA scan ($100-200)",
        "",
        "Competitors use 1-2 sources",
        "We use 6 = Defensible Moat"
    ]
    
    for line in moat_content:
        p = moat_text.add_paragraph()
        p.text = line
        if "Competitors" in line or "Defensible" in line:
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = colors['gold']
        else:
            p.font.size = Pt(15)
            p.font.color.rgb = colors['white']
        p.alignment = PP_ALIGN.LEFT
    
    # Right: UAE Advantages
    uae_box = slide.shapes.add_shape(1, Inches(5.1), Inches(col_y), Inches(4.4), Inches(col_height))
    uae_box.fill.solid()
    uae_box.fill.fore_color.rgb = RGBColor(40, 50, 70)
    uae_box.line.color.rgb = colors['gold']
    uae_box.line.width = Pt(2)
    
    uae_text = uae_box.text_frame
    uae_text.word_wrap = True
    uae_text.margin_left = Inches(0.2)
    uae_text.margin_right = Inches(0.2)
    uae_text.margin_top = Inches(0.1)
    
    pu = uae_text.paragraphs[0]
    pu.text = "🇦🇪 UAE-SPECIFIC ADVANTAGES"
    pu.font.size = Pt(24)
    pu.font.bold = True
    pu.font.color.rgb = colors['gold']
    pu.alignment = PP_ALIGN.CENTER
    
    uae_content = [
        "",
        "🧬 Emirati Genetics:",
        "   800K+ UAE genomes",
        "   30-40% higher accuracy",
        "",
        "🌙 Ramadan Protocols:",
        "   Fasting-aware AI",
        "   40%+ higher adherence",
        "",
        "🗣️ Arabic Language:",
        "   Native dialect, not translations",
        "",
        "Western competitors can't replicate"
    ]
    
    for line in uae_content:
        p = uae_text.add_paragraph()
        p.text = line
        if "Western" in line:
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = colors['gold']
        else:
            p.font.size = Pt(15)
            p.font.color.rgb = colors['white']
        p.alignment = PP_ALIGN.LEFT
    
    return slide


def create_slide3(prs, colors):
    """Slide 3: Seed → Series A → Dominance"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = colors['navy']
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = "From Seed to Market Leader: 18-Month Roadmap"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = colors['white']
    title_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(9), Inches(0.4))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "$15M Seed → $150M Series A → $1-3B Exit"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.color.rgb = colors['gold']
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Use of Funds (left)
    funds_box = slide.shapes.add_shape(1, Inches(0.5), Inches(1.5), Inches(4.5), Inches(2.2))
    funds_box.fill.solid()
    funds_box.fill.fore_color.rgb = RGBColor(40, 50, 70)
    funds_box.line.color.rgb = colors['green']
    funds_box.line.width = Pt(2)
    
    funds_text = funds_box.text_frame
    funds_text.word_wrap = True
    funds_text.margin_left = Inches(0.2)
    funds_text.margin_right = Inches(0.2)
    funds_text.margin_top = Inches(0.1)
    
    pf = funds_text.paragraphs[0]
    pf.text = "💰 USE OF $15M SEED FUNDS"
    pf.font.size = Pt(24)
    pf.font.bold = True
    pf.font.color.rgb = colors['green']
    pf.alignment = PP_ALIGN.CENTER
    
    funds_content = [
        "",
        "Team Building: $4M (27%)",
        "  CMO, CTO, VP Partnerships",
        "  8 Engineers, 3 Clinical, 2 Ops",
        "",
        "Product Development: $6M (40%)",
        "  Digital Twin AI engine: $2.5M",
        "  Mobile apps (iOS/Android): $1.5M",
        "  Backend + Thiqa integration: $1.5M",
        "",
        "Clinical Validation: $3M (20%)",
        "  500-person pilot (6 months)",
        "  Clinical study + IRB approval",
        "",
        "Regulatory & Operations: $2M (13%)"
    ]
    
    for line in funds_content:
        p = funds_text.add_paragraph()
        p.text = line
        if "$" in line and ":" in line:
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = colors['gold']
        else:
            p.font.size = Pt(14)
            p.font.color.rgb = colors['white']
        p.alignment = PP_ALIGN.LEFT
    
    # Timeline (right top)
    timeline_box = slide.shapes.add_shape(1, Inches(5.2), Inches(1.5), Inches(4.3), Inches(2.2))
    timeline_box.fill.solid()
    timeline_box.fill.fore_color.rgb = RGBColor(40, 50, 70)
    timeline_box.line.color.rgb = colors['blue']
    timeline_box.line.width = Pt(2)
    
    timeline_text = timeline_box.text_frame
    timeline_text.word_wrap = True
    timeline_text.margin_left = Inches(0.2)
    timeline_text.margin_right = Inches(0.2)
    timeline_text.margin_top = Inches(0.1)
    
    pt = timeline_text.paragraphs[0]
    pt.text = "📅 18-MONTH TIMELINE"
    pt.font.size = Pt(24)
    pt.font.bold = True
    pt.font.color.rgb = colors['blue']
    pt.alignment = PP_ALIGN.CENTER
    
    timeline_content = [
        "",
        "Months 0-6: BUILD",
        "  Hire team, build MVP",
        "  Thiqa MOU signed",
        "",
        "Months 6-12: VALIDATE",
        "  500-person pilot launch",
        "  80% retention @ Month 3",
        "  +10 Health Score average",
        "",
        "Months 12-18: SCALE",
        "  Publish research paper",
        "  Thiqa LOI (10,000 members)",
        "  $500K ARR (pilot converts)",
        "  Ready for Series A"
    ]
    
    for line in timeline_content:
        p = timeline_text.add_paragraph()
        p.text = line
        if "BUILD" in line or "VALIDATE" in line or "SCALE" in line:
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = colors['gold']
        else:
            p.font.size = Pt(14)
            p.font.color.rgb = colors['white']
        p.alignment = PP_ALIGN.LEFT
    
    # Bottom: Series A + Exit
    bottom_box = slide.shapes.add_shape(1, Inches(0.5), Inches(3.9), Inches(9), Inches(1.4))
    bottom_box.fill.solid()
    bottom_box.fill.fore_color.rgb = RGBColor(40, 50, 70)
    bottom_box.line.color.rgb = colors['gold']
    bottom_box.line.width = Pt(3)
    
    bottom_text = bottom_box.text_frame
    bottom_text.word_wrap = True
    bottom_text.margin_left = Inches(0.3)
    bottom_text.margin_right = Inches(0.3)
    bottom_text.margin_top = Inches(0.1)
    
    pb = bottom_text.paragraphs[0]
    pb.text = "🚀 SERIES A (MONTH 18) → 5-YEAR EXIT"
    pb.font.size = Pt(28)
    pb.font.bold = True
    pb.font.color.rgb = colors['gold']
    pb.alignment = PP_ALIGN.CENTER
    
    bottom_content = [
        "",
        "Series A Metrics: 500 users | 75% retention | $500K ARR | Thiqa LOI | Published research → $150M valuation (3x seed)",
        "",
        "5-Year Vision: 500K Daman/Hayat members | $90M ARR (Year 2) | 2M GCC members | $360M ARR (Year 4)",
        "",
        "Exit: $1-3B (Acquisition/IPO/Government) → 20-60x Returns for Seed Investors"
    ]
    
    for line in bottom_content:
        p = bottom_text.add_paragraph()
        p.text = line
        if "20-60x" in line or "$1-3B" in line:
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = colors['gold']
        else:
            p.font.size = Pt(16)
            p.font.color.rgb = colors['white']
        p.alignment = PP_ALIGN.CENTER
    
    return slide


def create_slide4(prs, colors):
    """Slide 4: UAE Healthspan Crisis Graph"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = colors['navy']
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = "The UAE Healthspan Crisis: 16 Years Lost"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = colors['white']
    title_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(9), Inches(0.4))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "7th Worst Globally and Getting Worse"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.color.rgb = colors['light_gray']
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Create line chart
    chart_data = CategoryChartData()
    chart_data.categories = ['2000', '2005', '2010', '2015', '2020', '2025']
    chart_data.add_series('UAE', (11.0, 12.5, 13.8, 14.9, 15.6, 16.0))
    chart_data.add_series('Global Average', (9.2, 9.1, 9.0, 8.9, 8.8, 8.7))
    chart_data.add_series('Top Performers', (6.5, 6.2, 5.9, 5.6, 5.3, 5.0))
    
    x, y, cx, cy = Inches(0.5), Inches(1.5), Inches(6.5), Inches(3.5)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
    ).chart
    
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    
    # Customize series colors
    # UAE - Red
    series_uae = chart.series[0]
    series_uae.format.line.color.rgb = colors['red']
    series_uae.format.line.width = Pt(4)
    
    # Global - Gray
    series_global = chart.series[1]
    series_global.format.line.color.rgb = RGBColor(128, 128, 128)
    series_global.format.line.width = Pt(2)
    series_global.format.line.dash_style = 2  # Dashed
    
    # Top Performers - Green
    series_top = chart.series[2]
    series_top.format.line.color.rgb = colors['green']
    series_top.format.line.width = Pt(2)
    
    # Statistics panel (right)
    stats_box = slide.shapes.add_shape(1, Inches(7.2), Inches(1.5), Inches(2.3), Inches(3.5))
    stats_box.fill.solid()
    stats_box.fill.fore_color.rgb = RGBColor(40, 50, 70)
    stats_box.line.color.rgb = colors['red']
    stats_box.line.width = Pt(2)
    
    stats_text = stats_box.text_frame
    stats_text.word_wrap = True
    stats_text.margin_left = Inches(0.15)
    stats_text.margin_right = Inches(0.15)
    stats_text.margin_top = Inches(0.1)
    
    ps = stats_text.paragraphs[0]
    ps.text = "📊 UAE CRISIS"
    ps.font.size = Pt(22)
    ps.font.bold = True
    ps.font.color.rgb = colors['red']
    ps.alignment = PP_ALIGN.CENTER
    
    stats_content = [
        "",
        "📉 16 years",
        "   in poor health",
        "",
        "🌍 7th worst",
        "   globally",
        "",
        "📈 +45% worse",
        "   since 2000",
        "",
        "🌎 1.8x global avg",
        "🏆 3.2x top performers",
        "",
        "👥 500K+ affected",
        "💰 $11B burden",
        "🏥 60%+ preventable"
    ]
    
    for line in stats_content:
        p = stats_text.add_paragraph()
        p.text = line
        if "16 years" in line or "7th worst" in line or "+45%" in line:
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = colors['gold']
        else:
            p.font.size = Pt(14)
            p.font.color.rgb = colors['white']
        p.alignment = PP_ALIGN.LEFT
    
    # Bottom text
    bottom_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.1), Inches(9), Inches(0.4))
    bottom_frame = bottom_box.text_frame
    bottom_frame.text = "While top countries improved 23%, UAE worsened 45%. Digital Twin AI can reverse this trajectory."
    bottom_para = bottom_frame.paragraphs[0]
    bottom_para.font.size = Pt(18)
    bottom_para.font.color.rgb = colors['white']
    bottom_para.alignment = PP_ALIGN.CENTER
    
    return slide


def create_slide5(prs, colors):
    """Slide 5: Digital Twin Market Validation"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = colors['navy']
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "Digital Twin Health: Proven by Academia, Funded by VCs"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(42)
    title_para.font.bold = True
    title_para.font.color.rgb = colors['white']
    title_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.75), Inches(9), Inches(0.35))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "$950M Valuation Proves Investors Believe"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(26)
    subtitle_para.font.color.rgb = colors['gold']
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Three columns: Academic, Twin Health, PrediSurge
    col_width = 2.9
    col_height = 2.6
    y_start = 1.2
    
    # Column 1: Academic (Mount Sinai)
    col1_x = 0.5
    col1_box = slide.shapes.add_shape(1, Inches(col1_x), Inches(y_start), Inches(col_width), Inches(col_height))
    col1_box.fill.solid()
    col1_box.fill.fore_color.rgb = RGBColor(40, 50, 70)
    col1_box.line.color.rgb = colors['blue_gray']
    col1_box.line.width = Pt(2)
    
    col1_text = col1_box.text_frame
    col1_text.word_wrap = True
    col1_text.margin_left = Inches(0.15)
    col1_text.margin_right = Inches(0.15)
    col1_text.margin_top = Inches(0.1)
    
    p1 = col1_text.paragraphs[0]
    p1.text = "🏥 ACADEMIC PROOF"
    p1.font.size = Pt(22)
    p1.font.bold = True
    p1.font.color.rgb = colors['blue_gray']
    p1.alignment = PP_ALIGN.CENTER
    
    content1 = [
        "",
        "Mount Sinai DigiTwin",
        "",
        "12M patient database",
        "",
        "Nature: 373 citations",
        "Lancet: 9 citations",
        "",
        "Dr. Zahi Fayad",
        "(TIME Magazine 2025)",
        "",
        "Proven: Diabetes, Sepsis,",
        "MS, Cancer, Pediatrics",
        "",
        "digitwinhealth.com"
    ]
    
    for line in content1:
        p = col1_text.add_paragraph()
        p.text = line
        if "12M" in line or "373" in line:
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = colors['gold']
        elif "Mount Sinai" in line:
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = colors['white']
        else:
            p.font.size = Pt(14)
            p.font.color.rgb = colors['white']
        p.alignment = PP_ALIGN.LEFT
    
    # Column 2: Twin Health
    col2_x = 3.55
    col2_box = slide.shapes.add_shape(1, Inches(col2_x), Inches(y_start), Inches(col_width), Inches(col_height))
    col2_box.fill.solid()
    col2_box.fill.fore_color.rgb = RGBColor(40, 50, 70)
    col2_box.line.color.rgb = colors['green']
    col2_box.line.width = Pt(3)
    
    col2_text = col2_box.text_frame
    col2_text.word_wrap = True
    col2_text.margin_left = Inches(0.15)
    col2_text.margin_right = Inches(0.15)
    col2_text.margin_top = Inches(0.1)
    
    p2 = col2_text.paragraphs[0]
    p2.text = "💰 TWIN HEALTH"
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = colors['green']
    p2.alignment = PP_ALIGN.CENTER
    
    content2 = [
        "",
        "Metabolic Digital Twin",
        "",
        "$950M Valuation",
        "(Aug 2025)",
        "",
        "$296M+ Total Raised",
        "Series E: $53M",
        "",
        "71% Diabetes Reversal",
        "(Cleveland Clinic)",
        "",
        "Investors: Iconiq,",
        "Temasek, Maj Invest",
        "",
        "$50M from Unicorn"
    ]
    
    for line in content2:
        p = col2_text.add_paragraph()
        p.text = line
        if "$950M" in line or "71%" in line:
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = colors['gold']
        elif "$296M" in line or "Unicorn" in line:
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = colors['gold']
        else:
            p.font.size = Pt(14)
            p.font.color.rgb = colors['white']
        p.alignment = PP_ALIGN.LEFT
    
    # Column 3: PrediSurge
    col3_x = 6.6
    col3_box = slide.shapes.add_shape(1, Inches(col3_x), Inches(y_start), Inches(col_width), Inches(col_height))
    col3_box.fill.solid()
    col3_box.fill.fore_color.rgb = RGBColor(40, 50, 70)
    col3_box.line.color.rgb = RGBColor(255, 99, 71)  # Red-orange
    col3_box.line.width = Pt(2)
    
    col3_text = col3_box.text_frame
    col3_text.word_wrap = True
    col3_text.margin_left = Inches(0.15)
    col3_text.margin_right = Inches(0.15)
    col3_text.margin_top = Inches(0.1)
    
    p3 = col3_text.paragraphs[0]
    p3.text = "🫀 PREDISURGE"
    p3.font.size = Pt(22)
    p3.font.bold = True
    p3.font.color.rgb = RGBColor(255, 99, 71)
    p3.alignment = PP_ALIGN.CENTER
    
    content3 = [
        "",
        "Cardiovascular",
        "Digital Twin",
        "",
        "€6.5M Series A",
        "($7M USD)",
        "(Oct 2023)",
        "",
        "3D Patient Models",
        "for Heart Surgery",
        "",
        "Reduces surgical risk",
        "Improves outcomes",
        "",
        "European leader"
    ]
    
    for line in content3:
        p = col3_text.add_paragraph()
        p.text = line
        if "€6.5M" in line or "$7M" in line:
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = colors['gold']
        else:
            p.font.size = Pt(14)
            p.font.color.rgb = colors['white']
        p.alignment = PP_ALIGN.LEFT
    
    # Bottom panels
    bottom_y = 4.0
    
    # Market Growth (left)
    market_box = slide.shapes.add_shape(1, Inches(0.5), Inches(bottom_y), Inches(4.5), Inches(1.3))
    market_box.fill.solid()
    market_box.fill.fore_color.rgb = RGBColor(40, 50, 70)
    market_box.line.color.rgb = colors['gold']
    market_box.line.width = Pt(2)
    
    market_text = market_box.text_frame
    market_text.word_wrap = True
    market_text.margin_left = Inches(0.2)
    market_text.margin_top = Inches(0.1)
    
    pm = market_text.paragraphs[0]
    pm.text = "📈 MARKET: $1.2B (2024) → $6.1B (2030) | 28.9% CAGR"
    pm.font.size = Pt(18)
    pm.font.bold = True
    pm.font.color.rgb = colors['gold']
    pm.alignment = PP_ALIGN.CENTER
    
    market_content = [
        "",
        "$300M+ invested in Digital Twin health (past 3 years)",
        "Near-unicorn valuations prove investor confidence"
    ]
    
    for line in market_content:
        p = market_text.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        p.font.color.rgb = colors['white']
        p.alignment = PP_ALIGN.CENTER
    
    # Hayat AI Positioning (right)
    hayat_box = slide.shapes.add_shape(1, Inches(5.2), Inches(bottom_y), Inches(4.3), Inches(1.3))
    hayat_box.fill.solid()
    hayat_box.fill.fore_color.rgb = RGBColor(40, 50, 70)
    hayat_box.line.color.rgb = colors['green']
    hayat_box.line.width = Pt(3)
    
    hayat_text = hayat_box.text_frame
    hayat_text.word_wrap = True
    hayat_text.margin_left = Inches(0.2)
    hayat_text.margin_top = Inches(0.1)
    
    ph = hayat_text.paragraphs[0]
    ph.text = "🇦🇪 HAYAT AI: Proven Tech + UAE Optimization"
    ph.font.size = Pt(18)
    ph.font.bold = True
    ph.font.color.rgb = colors['green']
    ph.alignment = PP_ALIGN.CENTER
    
    hayat_content = [
        "",
        "Twin Health (metabolic only) = $950M",
        "Hayat AI (CVD + Diabetes + Cancer) = Larger market",
        "Western tech + UAE moat = Defensible advantage"
    ]
    
    for line in hayat_content:
        p = hayat_text.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        if "Larger" in line or "Defensible" in line:
            p.font.bold = True
            p.font.color.rgb = colors['gold']
        else:
            p.font.color.rgb = colors['white']
        p.alignment = PP_ALIGN.CENTER
    
    return slide


if __name__ == "__main__":
    print("=" * 60)
    print("Creating Hayat AI 5 Dynamite Slides Presentation")
    print("=" * 60)
    print()
    
    output_file = create_presentation()
    
    print()
    print("=" * 60)
    print("✅ PRESENTATION COMPLETE!")
    print("=" * 60)
    print(f"File: {output_file}")
    print()
    print("Slides created:")
    print("1. The $11 Billion Crisis")
    print("2. Why Digital Twin AI Wins in the UAE")
    print("3. From Seed to Market Leader: 18-Month Roadmap")
    print("4. The UAE Healthspan Crisis: 16 Years Lost")
    print("5. Digital Twin Health: Proven by Academia, Funded by VCs")
    print()
