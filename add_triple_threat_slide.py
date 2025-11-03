#!/usr/bin/env python3
"""
Add Triple Threat slide to Daman-Hayat AI PowerPoint presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# File paths
INPUT_FILE = "/home/ubuntu/Daman-Hayat-AI-Project/presentation/Daman-Hayat_AI_Opportunity.pptx"
OUTPUT_FILE = "/home/ubuntu/Daman-Hayat-AI-Project/presentation/Daman-Hayat_AI_Opportunity_TripleThreat.pptx"

# Color palette
COLOR_RED = RGBColor(193, 18, 31)        # CVD
COLOR_ORANGE = RGBColor(255, 165, 0)     # Diabetes
COLOR_PURPLE = RGBColor(114, 9, 183)     # Cancer
COLOR_TEAL = RGBColor(5, 102, 141)       # Positive/Solution
COLOR_DARK_GRAY = RGBColor(44, 62, 80)   # Body text
COLOR_LIGHT_GRAY = RGBColor(240, 240, 240) # Background

def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=COLOR_DARK_GRAY, align=PP_ALIGN.CENTER):
    """Helper function to add text box with formatting"""
    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    text_frame = textbox.text_frame
    text_frame.text = text
    text_frame.word_wrap = True
    
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = align
    
    run = paragraph.runs[0]
    run.font.name = "Noto Sans"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    
    return textbox

def create_triple_threat_slide(prs):
    """Create the Triple Threat slide"""
    
    # Add blank slide (using layout 6 - blank)
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Set background color to white
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Add title
    add_text_box(slide, 0.5, 0.3, 9, 0.8, 
                 "The Triple Threat: UAE's Preventable Health Crisis",
                 font_size=40, bold=True, color=COLOR_TEAL)
    
    # Column dimensions
    col_width = 2.8
    col_height = 5.0
    col_top = 1.5
    col_spacing = 0.2
    
    # Column 1: CVD (left)
    col1_left = 0.5
    
    # CVD Header
    add_text_box(slide, col1_left, col_top, col_width, 0.4,
                 "🔴 CVD", font_size=28, bold=True, color=COLOR_RED)
    add_text_box(slide, col1_left, col_top + 0.4, col_width, 0.3,
                 "#1 KILLER", font_size=16, bold=True, color=COLOR_RED)
    
    # CVD Statistics
    y_offset = col_top + 0.9
    
    add_text_box(slide, col1_left, y_offset, col_width, 0.5,
                 "40%", font_size=48, bold=True, color=COLOR_RED)
    add_text_box(slide, col1_left, y_offset + 0.5, col_width, 0.3,
                 "of all deaths", font_size=14, color=COLOR_DARK_GRAY)
    
    y_offset += 1.0
    add_text_box(slide, col1_left, y_offset, col_width, 0.5,
                 "2x", font_size=42, bold=True, color=COLOR_RED)
    add_text_box(slide, col1_left, y_offset + 0.5, col_width, 0.3,
                 "cancer rate", font_size=14, color=COLOR_DARK_GRAY)
    
    y_offset += 1.0
    add_text_box(slide, col1_left, y_offset, col_width, 0.5,
                 "314/100K", font_size=36, bold=True, color=COLOR_RED)
    add_text_box(slide, col1_left, y_offset + 0.5, col_width, 0.3,
                 "death rate", font_size=14, color=COLOR_DARK_GRAY)
    
    y_offset += 1.0
    add_text_box(slide, col1_left, y_offset, col_width, 0.5,
                 "80%", font_size=42, bold=True, color=COLOR_TEAL)
    add_text_box(slide, col1_left, y_offset + 0.5, col_width, 0.3,
                 "preventable", font_size=14, color=COLOR_TEAL)
    
    # CVD Detection methods
    add_text_box(slide, col1_left, col_top + col_height - 0.5, col_width, 0.5,
                 "Detectable: Retinal AI, DEXA, Wearables", 
                 font_size=11, color=COLOR_DARK_GRAY)
    
    # Column 2: Diabetes (center)
    col2_left = col1_left + col_width + col_spacing
    
    # Diabetes Header
    add_text_box(slide, col2_left, col_top, col_width, 0.4,
                 "🟡 DIABETES", font_size=28, bold=True, color=COLOR_ORANGE)
    add_text_box(slide, col2_left, col_top + 0.4, col_width, 0.3,
                 "EPIDEMIC", font_size=16, bold=True, color=COLOR_ORANGE)
    
    # Diabetes Statistics
    y_offset = col_top + 0.9
    
    add_text_box(slide, col2_left, y_offset, col_width, 0.5,
                 "20.7%", font_size=48, bold=True, color=COLOR_ORANGE)
    add_text_box(slide, col2_left, y_offset + 0.5, col_width, 0.3,
                 "prevalence", font_size=14, color=COLOR_DARK_GRAY)
    
    y_offset += 1.0
    add_text_box(slide, col2_left, y_offset, col_width, 0.5,
                 "2.5x", font_size=42, bold=True, color=COLOR_ORANGE)
    add_text_box(slide, col2_left, y_offset + 0.5, col_width, 0.3,
                 "global average", font_size=14, color=COLOR_DARK_GRAY)
    
    y_offset += 1.0
    add_text_box(slide, col2_left, y_offset, col_width, 0.5,
                 "1.27M", font_size=36, bold=True, color=COLOR_ORANGE)
    add_text_box(slide, col2_left, y_offset + 0.5, col_width, 0.3,
                 "total cases", font_size=14, color=COLOR_DARK_GRAY)
    
    y_offset += 1.0
    add_text_box(slide, col2_left, y_offset, col_width, 0.5,
                 "90%", font_size=42, bold=True, color=COLOR_TEAL)
    add_text_box(slide, col2_left, y_offset + 0.5, col_width, 0.3,
                 "Type 2 (preventable)", font_size=14, color=COLOR_TEAL)
    
    # Diabetes Detection methods
    add_text_box(slide, col2_left, col_top + col_height - 0.5, col_width, 0.5,
                 "Detectable: Blood, Genome, DEXA", 
                 font_size=11, color=COLOR_DARK_GRAY)
    
    # Column 3: Cancer (right)
    col3_left = col2_left + col_width + col_spacing
    
    # Cancer Header
    add_text_box(slide, col3_left, col_top, col_width, 0.4,
                 "🔵 CANCER", font_size=28, bold=True, color=COLOR_PURPLE)
    add_text_box(slide, col3_left, col_top + 0.4, col_width, 0.3,
                 "RISING THREAT", font_size=16, bold=True, color=COLOR_PURPLE)
    
    # Cancer Statistics
    y_offset = col_top + 0.9
    
    add_text_box(slide, col3_left, y_offset, col_width, 0.5,
                 "12%", font_size=48, bold=True, color=COLOR_PURPLE)
    add_text_box(slide, col3_left, y_offset + 0.5, col_width, 0.3,
                 "of all deaths", font_size=14, color=COLOR_DARK_GRAY)
    
    y_offset += 1.0
    add_text_box(slide, col3_left, y_offset, col_width, 0.5,
                 "15.4%", font_size=42, bold=True, color=COLOR_PURPLE)
    add_text_box(slide, col3_left, y_offset + 0.5, col_width, 0.3,
                 "young-onset", font_size=14, color=COLOR_DARK_GRAY)
    
    y_offset += 1.0
    add_text_box(slide, col3_left, y_offset, col_width, 0.5,
                 "5,526", font_size=36, bold=True, color=COLOR_PURPLE)
    add_text_box(slide, col3_left, y_offset + 0.5, col_width, 0.3,
                 "new cases/year", font_size=14, color=COLOR_DARK_GRAY)
    
    y_offset += 1.0
    add_text_box(slide, col3_left, y_offset, col_width, 0.5,
                 "10 years", font_size=36, bold=True, color=COLOR_TEAL)
    add_text_box(slide, col3_left, y_offset + 0.5, col_width, 0.3,
                 "early detection", font_size=14, color=COLOR_TEAL)
    
    # Cancer Detection methods
    add_text_box(slide, col3_left, col_top + col_height - 0.5, col_width, 0.5,
                 "Detectable: AI Liquid Biopsy, Genome", 
                 font_size=11, color=COLOR_DARK_GRAY)
    
    # Bottom section: Combined Impact
    bottom_top = col_top + col_height + 0.3
    
    # Add background rectangle
    left = Inches(0.5)
    top = Inches(bottom_top)
    width = Inches(9)
    height = Inches(0.8)
    shape = slide.shapes.add_shape(1, left, top, width, height)  # 1 = rectangle
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_LIGHT_GRAY
    shape.line.color.rgb = COLOR_TEAL
    shape.line.width = Pt(2)
    
    # Combined impact text
    add_text_box(slide, 0.5, bottom_top + 0.1, 9, 0.3,
                 "COMBINED IMPACT: 52% of all UAE deaths + $11B annual cost",
                 font_size=20, bold=True, color=COLOR_RED)
    
    add_text_box(slide, 0.5, bottom_top + 0.45, 9, 0.3,
                 "✅ Hayat AI detects and prevents all three simultaneously",
                 font_size=18, bold=True, color=COLOR_TEAL)
    
    # Sources footer
    add_text_box(slide, 0.5, 7.2, 9, 0.2,
                 "Sources: American Heart Association (2024), IDF (2024), WHO IARC (2024), UNDP (2024), PMC (2023-2024)",
                 font_size=9, color=COLOR_DARK_GRAY)
    
    return slide

def main():
    """Main function to add Triple Threat slide"""
    print("Loading presentation...")
    prs = Presentation(INPUT_FILE)
    
    print(f"Original presentation has {len(prs.slides)} slides")
    
    print("Creating Triple Threat slide...")
    create_triple_threat_slide(prs)
    
    print(f"Updated presentation has {len(prs.slides)} slides")
    
    print(f"Saving to {OUTPUT_FILE}...")
    prs.save(OUTPUT_FILE)
    
    print("✅ Triple Threat slide added successfully!")
    print(f"📁 Output file: {OUTPUT_FILE}")

if __name__ == "__main__":
    main()
