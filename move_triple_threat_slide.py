#!/usr/bin/env python3
"""
Move Triple Threat slide from position #22 to position #5
Creates logical flow: Problem → Root Causes → Solution
"""

from pptx import Presentation
from collections.abc import Sequence

# File paths
INPUT_FILE = "/home/ubuntu/Daman-Hayat-AI-Project/presentation/Daman-Hayat_AI_Opportunity_Updated.pptx"
OUTPUT_FILE = "/home/ubuntu/Daman-Hayat-AI-Project/presentation/Daman-Hayat_AI_Opportunity_Final.pptx"

def move_slide(prs, old_index, new_index):
    """
    Move a slide from old_index to new_index
    Note: Indices are 0-based
    """
    # Get the XML of all slides
    slides = list(prs.slides._sldIdLst)
    
    # Remove the slide from old position
    slide_to_move = slides.pop(old_index)
    
    # Insert at new position
    slides.insert(new_index, slide_to_move)
    
    # Update the XML
    prs.slides._sldIdLst[:] = slides
    
    return prs

def main():
    """Main function to move Triple Threat slide"""
    print("Loading presentation...")
    prs = Presentation(INPUT_FILE)
    
    total_slides = len(prs.slides)
    print(f"Presentation has {total_slides} slides")
    
    # Move from position 22 (index 21) to position 5 (index 4)
    old_position = 22  # Human-readable position
    new_position = 5   # Human-readable position
    
    old_index = old_position - 1  # Convert to 0-based index
    new_index = new_position - 1  # Convert to 0-based index
    
    print(f"\nMoving slide from position {old_position} to position {new_position}...")
    print(f"(Index {old_index} → Index {new_index})")
    
    # Move the slide
    move_slide(prs, old_index, new_index)
    
    print(f"\nSaving to {OUTPUT_FILE}...")
    prs.save(OUTPUT_FILE)
    
    print("\n✅ Triple Threat slide moved successfully!")
    print(f"📁 Output file: {OUTPUT_FILE}")
    print(f"\n📊 New slide order:")
    print(f"  Slide 1-4: Introduction and Problem")
    print(f"  Slide 5: Triple Threat (ROOT CAUSES) ← MOVED HERE")
    print(f"  Slide 6+: Solution and Details")
    print(f"\n🎯 Narrative flow: Problem → Root Causes → Solution")

if __name__ == "__main__":
    main()
