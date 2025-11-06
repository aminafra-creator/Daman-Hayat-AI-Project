#!/usr/bin/env python3
"""
Add Digital Twin Technology slide to presentation
Position: After Triple Threat slide (position #6)
Layout: 3-column design with UAE customization focus
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# File paths
INPUT_FILE = "/home/ubuntu/Daman-Hayat-AI-Project/presentation/Daman-Hayat_AI_Opportunity_Final.pptx"
OUTPUT_FILE = "/home/ubuntu/Daman-Hayat-AI-Project/presentation/Daman-Hayat_AI_Opportunity_DigitalTwin.pptx"

# Color palette
COLOR_TEAL = RGBColor(5, 102, 141)
COLOR_DARK_GRAY = RGBColor(44, 62, 80)
COLOR_LIGHT_GRAY = RGBColor(240, 240, 240)
COLOR_GREEN = RGBColor(0, 150, 100)
COLOR_BLUE = RGBColor(0, 120, 215)
COLOR_ORANGE = RGBColor(255, 140, 0)

def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=COLOR_DARK_GRAY, align=PP_ALIGN.LEFT):
    """Helper function to add text box with formatting"""
    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    text_frame = textbox.text_frame
    text_frame.text = text
    text_frame.word_wrap = True
    
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = align
    
    run = paragraph.runs[0]
    run.font.name = "Noto Sans"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    
    return textbox

def create_digital_twin_slide(prs):
    """Create the Digital Twin Technology slide"""
    
    # Add blank slide
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Set background color to white
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Add title
    add_text_box(slide, 0.5, 0.3, 9, 0.6, 
                 "Digital Twin Technology—Customized for the UAE",
                 font_size=36, bold=True, color=COLOR_TEAL, align=PP_ALIGN.CENTER)
    
    # Add subtitle
    add_text_box(slide, 0.5, 0.9, 9, 0.3,
                 "Global Innovation, UAE Customization",
                 font_size=20, bold=False, color=COLOR_DARK_GRAY, align=PP_ALIGN.CENTER)
    
    # Column dimensions
    col_width = 2.8
    col_height = 4.5
    col_top = 1.5
    col_spacing = 0.2
    
    # Column 1: Proven Concept (left)
    col1_left = 0.5
    
    # Add icon/emoji
    add_text_box(slide, col1_left, col_top, col_width, 0.4,
                 "🌍", font_size=32, align=PP_ALIGN.CENTER)
    
    # Column 1 Header
    add_text_box(slide, col1_left, col_top + 0.5, col_width, 0.4,
                 "PROVEN CONCEPT", font_size=20, bold=True, color=COLOR_GREEN, align=PP_ALIGN.CENTER)
    
    # Column 1 Content
    y_offset = col_top + 1.0
    
    add_text_box(slide, col1_left + 0.1, y_offset, col_width - 0.2, 0.8,
                 "Mount Sinai & leading institutions pioneering Digital Twin for health optimization",
                 font_size=14, color=COLOR_DARK_GRAY, align=PP_ALIGN.LEFT)
    
    y_offset += 1.0
    add_text_box(slide, col1_left + 0.1, y_offset, col_width - 0.2, 1.2,
                 "We're bringing this proven concept to the UAE—with critical customization",
                 font_size=14, bold=True, color=COLOR_TEAL, align=PP_ALIGN.LEFT)
    
    y_offset += 1.4
    add_text_box(slide, col1_left + 0.1, y_offset, col_width - 0.2, 0.8,
                 "Transforming healthcare from treating illness to improving wellness",
                 font_size=13, color=COLOR_DARK_GRAY, align=PP_ALIGN.LEFT)
    
    # Column 2: UAE-Specific (center)
    col2_left = col1_left + col_width + col_spacing
    
    # Add icon/emoji
    add_text_box(slide, col2_left, col_top, col_width, 0.4,
                 "🇦🇪", font_size=32, align=PP_ALIGN.CENTER)
    
    # Column 2 Header
    add_text_box(slide, col2_left, col_top + 0.5, col_width, 0.4,
                 "UAE-SPECIFIC", font_size=20, bold=True, color=COLOR_BLUE, align=PP_ALIGN.CENTER)
    
    # Column 2 Content - Bullet points
    y_offset = col_top + 1.0
    
    bullets = [
        "• Emirati Genome Database",
        "• Ramadan-aware protocols",
        "• Arabic interface",
        "• Halal nutrition guidance",
        "• Cultural adaptation"
    ]
    
    for bullet in bullets:
        add_text_box(slide, col2_left + 0.1, y_offset, col_width - 0.2, 0.35,
                     bullet, font_size=14, color=COLOR_DARK_GRAY, align=PP_ALIGN.LEFT)
        y_offset += 0.45
    
    y_offset += 0.3
    add_text_box(slide, col2_left + 0.1, y_offset, col_width - 0.2, 0.8,
                 "No Western platform can match this customization",
                 font_size=13, bold=True, color=COLOR_BLUE, align=PP_ALIGN.LEFT)
    
    # Column 3: Triple Threat (right)
    col3_left = col2_left + col_width + col_spacing
    
    # Add icon/emoji
    add_text_box(slide, col3_left, col_top, col_width, 0.4,
                 "🎯", font_size=32, align=PP_ALIGN.CENTER)
    
    # Column 3 Header
    add_text_box(slide, col3_left, col_top + 0.5, col_width, 0.4,
                 "TRIPLE THREAT FOCUS", font_size=20, bold=True, color=COLOR_ORANGE, align=PP_ALIGN.CENTER)
    
    # Column 3 Content
    y_offset = col_top + 1.0
    
    diseases = [
        ("• CVD", "(40% deaths)"),
        ("• Diabetes", "(20.7% prevalence)"),
        ("• Cancer", "(12% deaths)")
    ]
    
    for disease, stat in diseases:
        add_text_box(slide, col3_left + 0.1, y_offset, col_width - 0.2, 0.3,
                     disease, font_size=15, bold=True, color=COLOR_DARK_GRAY, align=PP_ALIGN.LEFT)
        add_text_box(slide, col3_left + 0.1, y_offset + 0.25, col_width - 0.2, 0.25,
                     stat, font_size=12, color=COLOR_DARK_GRAY, align=PP_ALIGN.LEFT)
        y_offset += 0.6
    
    y_offset += 0.2
    add_text_box(slide, col3_left + 0.1, y_offset, col_width - 0.2, 0.6,
                 "52% of deaths\n$11B annual cost",
                 font_size=14, bold=True, color=COLOR_ORANGE, align=PP_ALIGN.LEFT)
    
    y_offset += 0.8
    add_text_box(slide, col3_left + 0.1, y_offset, col_width - 0.2, 0.6,
                 "5-10 year early detection before symptoms",
                 font_size=13, color=COLOR_TEAL, align=PP_ALIGN.LEFT)
    
    # Bottom section
    bottom_top = col_top + col_height + 0.3
    
    # Add background rectangle
    left = Inches(0.5)
    top = Inches(bottom_top)
    width = Inches(9)
    height = Inches(0.6)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_LIGHT_GRAY
    shape.line.color.rgb = COLOR_TEAL
    shape.line.width = Pt(2)
    
    # Bottom text
    add_text_box(slide, 0.5, bottom_top + 0.15, 9, 0.3,
                 "Commercial Platform • Immediate Deployment • Proven PPP Model",
                 font_size=18, bold=True, color=COLOR_TEAL, align=PP_ALIGN.CENTER)
    
    # Sources footer
    add_text_box(slide, 0.5, 7.2, 9, 0.2,
                 "Sources: Mount Sinai DigiTwin Project, UAE Genome Database, IDF (2024), UNDP (2024)",
                 font_size=9, color=COLOR_DARK_GRAY, align=PP_ALIGN.CENTER)
    
    return slide

def insert_slide_at_position(prs, new_slide, position):
    """Insert a slide at a specific position (1-indexed)"""
    # Get XML of all slides
    slides = list(prs.slides._sldIdLst)
    
    # Remove the new slide from the end
    new_slide_xml = slides.pop()
    
    # Insert at desired position (convert to 0-indexed)
    slides.insert(position - 1, new_slide_xml)
    
    # Update the XML
    prs.slides._sldIdLst[:] = slides
    
    return prs

def main():
    """Main function to add Digital Twin slide"""
    print("Loading presentation...")
    prs = Presentation(INPUT_FILE)
    
    original_count = len(prs.slides)
    print(f"Original presentation has {original_count} slides")
    
    print("\nCreating Digital Twin Technology slide...")
    new_slide = create_digital_twin_slide(prs)
    
    print(f"Slide created, now has {len(prs.slides)} slides")
    
    # Insert at position 6 (after Triple Threat at position 5)
    target_position = 6
    print(f"\nMoving slide to position {target_position}...")
    insert_slide_at_position(prs, new_slide, target_position)
    
    print(f"\nSaving to {OUTPUT_FILE}...")
    prs.save(OUTPUT_FILE)
    
    print("\n✅ Digital Twin slide added successfully!")
    print(f"📁 Output file: {OUTPUT_FILE}")
    print(f"\n📊 New presentation structure:")
    print(f"  Slides 1-4: Introduction and Problem")
    print(f"  Slide 5: Triple Threat (Root Causes)")
    print(f"  Slide 6: Digital Twin Technology (Vision) ← NEW SLIDE")
    print(f"  Slides 7+: Solution Details")
    print(f"\n🎯 Narrative flow: Problem → Root Causes → Vision → Solution")

if __name__ == "__main__":
    main()
