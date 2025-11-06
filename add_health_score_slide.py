from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Load the latest presentation
prs = Presentation('/home/ubuntu/Daman-Hayat-AI-Project/presentation/Daman-Hayat_AI_Opportunity_DigitalTwin.pptx')

# Add a blank slide
blank_slide_layout = prs.slide_layouts[6]  # Blank layout
slide = prs.slides.add_slide(blank_slide_layout)

# Set slide background to dark
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(20, 30, 48)  # Dark blue background

# Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
title_frame = title_box.text_frame
title_frame.text = "HAYAT HEALTH SCORE: Measuring Your Healthspan"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(36)
title_para.font.bold = True
title_para.font.color.rgb = RGBColor(255, 255, 255)
title_para.alignment = PP_ALIGN.CENTER

# Subtitle
subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(0.4))
subtitle_frame = subtitle_box.text_frame
subtitle_frame.text = "One Comprehensive Score + Six Actionable Pillars (Blueprint-Optimized)"
subtitle_para = subtitle_frame.paragraphs[0]
subtitle_para.font.size = Pt(18)
subtitle_para.font.color.rgb = RGBColor(200, 200, 200)
subtitle_para.alignment = PP_ALIGN.CENTER

# Main score box (center)
main_score_box = slide.shapes.add_textbox(Inches(3.5), Inches(1.8), Inches(3), Inches(1.2))
main_score_frame = main_score_box.text_frame
main_score_frame.text = "HAYAT HEALTH SCORE\n0-100"
for para in main_score_frame.paragraphs:
    para.font.size = Pt(28)
    para.font.bold = True
    para.font.color.rgb = RGBColor(0, 180, 216)  # Cyan
    para.alignment = PP_ALIGN.CENTER

# Add border to main score box
shape = slide.shapes[-1]
line = shape.line
line.color.rgb = RGBColor(0, 180, 216)
line.width = Pt(3)

# Six pillars (3x2 grid)
pillars = [
    {"name": "😴 Sleep & Recovery", "weight": "25%", "color": RGBColor(114, 9, 183), "row": 0, "col": 0},  # Purple
    {"name": "🫀 Cardiovascular", "weight": "25%", "color": RGBColor(193, 18, 31), "row": 0, "col": 1},  # Red
    {"name": "🩸 Metabolic Health", "weight": "20%", "color": RGBColor(255, 140, 0), "row": 0, "col": 2},  # Orange
    {"name": "🧠 Cognitive Health", "weight": "10%", "color": RGBColor(5, 102, 141), "row": 1, "col": 0},  # Teal
    {"name": "😊 Mental Health", "weight": "10%", "color": RGBColor(0, 150, 100), "row": 1, "col": 1},  # Green
    {"name": "🦠 Microbiome", "weight": "10%", "color": RGBColor(100, 100, 100), "row": 1, "col": 2},  # Gray
]

# Starting position for pillars
start_x = 0.5
start_y = 3.5
box_width = 3.0
box_height = 1.2
spacing_x = 0.2
spacing_y = 0.3

for pillar in pillars:
    # Calculate position
    x = start_x + pillar["col"] * (box_width + spacing_x)
    y = start_y + pillar["row"] * (box_height + spacing_y)
    
    # Add pillar box
    pillar_box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(box_width), Inches(box_height))
    pillar_frame = pillar_box.text_frame
    pillar_frame.word_wrap = True
    
    # Pillar name
    p1 = pillar_frame.paragraphs[0]
    p1.text = pillar["name"]
    p1.font.size = Pt(18)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(255, 255, 255)
    p1.alignment = PP_ALIGN.CENTER
    
    # Weight
    p2 = pillar_frame.add_paragraph()
    p2.text = pillar["weight"]
    p2.font.size = Pt(32)
    p2.font.bold = True
    p2.font.color.rgb = pillar["color"]
    p2.alignment = PP_ALIGN.CENTER
    
    # Add border
    shape = slide.shapes[-1]
    line = shape.line
    line.color.rgb = pillar["color"]
    line.width = Pt(2)

# Bottom text
bottom_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.6))
bottom_frame = bottom_box.text_frame
bottom_frame.text = "Based on Bryan Johnson's Blueprint Protocol • Predicts Health 12-18 Months Ahead • UAE-Customized"
bottom_para = bottom_frame.paragraphs[0]
bottom_para.font.size = Pt(14)
bottom_para.font.color.rgb = RGBColor(150, 150, 150)
bottom_para.alignment = PP_ALIGN.CENTER

# Save presentation
output_path = '/home/ubuntu/Daman-Hayat-AI-Project/presentation/Daman-Hayat_AI_Opportunity_HealthScore.pptx'
prs.save(output_path)
print(f"✅ Health Score slide added successfully!")
print(f"📁 Saved to: {output_path}")
print(f"📊 Total slides: {len(prs.slides)}")
